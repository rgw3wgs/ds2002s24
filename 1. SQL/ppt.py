from pptx import Presentation

# Data for slides
slides_data = [
    {
        "title": "1a. Display Actor Names",
        "content": "SELECT first_name, last_name FROM actor;"
    },
    {
        "title": "1b. Actor Name in Upper Case",
        "content": "SELECT UPPER(CONCAT(first_name, ' ', last_name)) AS 'Actor Name' FROM actor;"
    },
    {
        "title": "2a. Search Actor by First Name 'Joe'",
        "content": "SELECT actor_id, first_name, last_name FROM actor WHERE first_name LIKE 'Joe';"
    },
    {
        "title": "2b. Actors with 'GEN' in Last Name",
        "content": "SELECT first_name, last_name FROM actor WHERE last_name LIKE '%GEN%';"
    },
    {
        "title": "2c. Actors with 'LI' in Last Name",
        "content": "SELECT first_name, last_name FROM actor WHERE last_name LIKE '%LI%' ORDER BY last_name, first_name;"
    },
    {
        "title": "2d. Specific Countries' IDs",
        "content": "SELECT country_id, country FROM country WHERE country IN ('Afghanistan', 'Bangladesh', 'China');"
    },
    {
        "title": "3a. Add middle_name Column",
        "content": "ALTER TABLE actor ADD COLUMN middle_name VARCHAR(30) AFTER first_name;"
    },
    {
        "title": "3b. Change Data Type of middle_name",
        "content": "ALTER TABLE actor MODIFY COLUMN middle_name BLOB;"
    },
    {
        "title": "3c. Delete middle_name Column",
        "content": "ALTER TABLE actor DROP COLUMN middle_name;"
    },
    {
        "title": "4a. List Last Names and Counts",
        "content": "SELECT last_name AS 'Last Name', COUNT(last_name) AS 'Last Name Count' FROM actor GROUP BY last_name;"
    },
    {
        "title": "4b. Last Names Shared by At Least Two Actors",
        "content": "SELECT last_name AS 'Last Name', COUNT(last_name) AS 'Last Name Count' FROM actor GROUP BY last_name HAVING COUNT(last_name) > 1;"
    },
    {
        "title": "4c. Fix HARPO WILLIAMS Record",
        "content": "UPDATE actor SET first_name = 'HARPO' WHERE first_name = 'Groucho' AND last_name = 'Williams';"
    },
    {
        "title": "4d. Correct Name from HARPO to GROUCHO",
        "content": "UPDATE actor SET first_name = CASE WHEN first_name = 'Harpo' THEN 'GROUCHO' WHEN first_name = 'Groucho' THEN 'MUCHO GROUCHO' ELSE first_name END;"
    },
    # ... (Add all other slides following the same structure)
]

# Create a presentation object
prs = Presentation()

for slide_data in slides_data:
    slide_layout = prs.slide_layouts[1]  # Using layout 1 for title and content
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]  # Index 1 is the content placeholder in this layout

    # Set title and content for each slide
    title.text = slide_data["title"]
    content.text = slide_data["content"]

# Save the presentation
prs.save('SQL_Presentation.pptx')

print("Presentation created successfully!")
